#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
ppt_generator/ppt_kit.py
========================

可复用的 PPT 绘图工具箱（主题 + 原语）。

设计意图：每次输入内容不同，PPT 布局也应不同。因此我们不提供"固定模板"，
而是由 Agent **根据 Markdown 内容每次生成一个贴合内容的 Python 脚本**；
该脚本 `import ppt_kit` 复用本文件提供的主题配色与绘图原语，只需关注
"把哪些内容画到什么位置"，而不必重复造轮子。

用法（在 Agent 生成的脚本中）：

    import sys
    SKILL_DIR = r"<skill 目录绝对路径>"   # 由 skill_tool 返回的 skill_directory
    sys.path.insert(0, SKILL_DIR)
    from ppt_kit import (
        new_presentation, add_header, add_divider, add_card, add_bullets,
        add_table, add_quote, set_bg, add_text, add_rect, save,
        COLOR_PRIMARY, COLOR_ACCENT, COLOR_WARN, COLOR_TEXT, COLOR_LIGHT,
        COLOR_BG, COLOR_WHITE, FONT_CN,
    )

    prs, slide = new_presentation()
    add_header(slide, "标题", "副标题", badge="DeepAgent")
    add_divider(slide)
    add_card(slide, left, top, w, h, "栏标题", ["要点1","要点2"], accent="primary")
    add_table(slide, left, top, w, h, headers=["A","B"], rows=[["1","2"]])
    add_quote(slide, "结论金句")
    save(prs, r"<output 绝对路径>/xxx.pptx")

首次 import 时自动 `pip install python-pptx`。
"""
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    import sys
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "--quiet", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ----------------------- 主题配色 -----------------------
COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)          # 浅灰背景
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 主色-深蓝
COLOR_ACCENT = RGBColor(0x10, 0xA3, 0x7A)       # 强调色-绿
COLOR_WARN = RGBColor(0xE8, 0x6E, 0x3C)         # 警示-橙
COLOR_TEXT = RGBColor(0x2C, 0x33, 0x45)         # 正文-深灰
COLOR_LIGHT = RGBColor(0x6B, 0x72, 0x80)        # 次要文字-中灰
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
COLOR_ALT_ROW = RGBColor(0xF0, 0xF4, 0xFA)      # 表格隔行

ACCENT_MAP = {
    "warn": COLOR_WARN,
    "primary": COLOR_PRIMARY,
    "accent": COLOR_ACCENT,
    "info": COLOR_PRIMARY,
    None: COLOR_PRIMARY,
}

FONT_CN = "Microsoft YaHei"


def accent_color(name=None):
    """按名称取强调色；未知名称回退主色。"""
    if name is None:
        return COLOR_PRIMARY
    return ACCENT_MAP.get(str(name).strip().lower(), COLOR_PRIMARY)


# ----------------------- 基础原语 -----------------------
def new_presentation():
    """新建 16:9 宽屏演示，返回 (prs, 首张空白页 slide)。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG)
    return prs, slide


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=None):
    """多行文本（按 \\n 分段）。"""
    if color is None:
        color = COLOR_TEXT
    if font_name is None:
        font_name = FONT_CN
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_bullets(slide, left, top, width, height, items, size=11,
                color=None, marker="•", space_after=4):
    """纯项目符号列表（不带卡片背景）。"""
    if color is None:
        color = COLOR_TEXT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = (marker + " " + item) if marker else item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT_CN
    return tb


# ----------------------- 复合组件 -----------------------
def add_header(slide, title, subtitle="", badge=""):
    """顶部标题区：主标题 + 副标题 + 右上角标识 + 分隔线。"""
    add_text(slide, Inches(0.6), Inches(0.32), Inches(9.5), Inches(0.7),
             title, size=30, bold=True, color=COLOR_PRIMARY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.98), Inches(9.5), Inches(0.4),
                 subtitle, size=15, bold=False, color=COLOR_LIGHT,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if badge:
        add_text(slide, Inches(10.0), Inches(0.4), Inches(3.1), Inches(0.5),
                 badge, size=11, bold=True, color=COLOR_ACCENT,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_divider(slide, top=Inches(1.42), left=Inches(0.6),
                width=Inches(12.13), color=None):
    add_rect(slide, left, top, width, Pt(2.5), color or COLOR_PRIMARY)


def add_card(slide, left, top, width, height, title, items,
             accent="primary", item_size=11):
    """带标题条的卡片 + 项目符号列表。"""
    accent_c = accent_color(accent)
    add_rect(slide, left, top, width, height, COLOR_CARD_BG, COLOR_DIVIDER)
    title_h = Inches(0.42)
    add_rect(slide, left, top, width, title_h, accent_c)
    add_text(slide, left, top, width, title_h, title, size=14, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    body_top = top + title_h + Inches(0.08)
    body_h = height - title_h - Inches(0.16)
    add_bullets(slide, left + Inches(0.12), body_top,
                width - Inches(0.24), body_h, items, size=item_size)


def add_table(slide, left, top, width, height, headers, rows,
              accent="primary", body_size=11, header_size=12,
              col_widths=None):
    """带主题样式的表格。

    - 表头填充强调色、白字加粗；
    - 正文白底，隔行用 COLOR_ALT_ROW 浅蓝；
    - `col_widths`（Emu 列表，可选）按比例分配列宽。
    """
    accent_c = accent_color(accent)
    nrows = len(rows) + 1
    ncols = len(headers)
    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = gf.table

    # 表头
    for j, h in enumerate(headers):
        _set_cell(table.cell(0, j), str(h), accent_c, COLOR_WHITE,
                  bold=True, size=header_size)
    # 正文
    for i, row in enumerate(rows):
        fill = COLOR_ALT_ROW if i % 2 == 1 else COLOR_WHITE
        for j, val in enumerate(row):
            _set_cell(table.cell(i + 1, j), str(val), fill, COLOR_TEXT,
                      size=body_size)
    # 列宽
    if col_widths:
        for j, w in enumerate(col_widths):
            if j < ncols:
                table.columns[j].width = w
    # 关闭默认表格样式带来的边框杂色：设置首行强调已足够
    return table


def _set_cell(cell, text, fill, color, bold=False, size=11,
              align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Emu(54000)
    cell.margin_right = Emu(54000)
    cell.margin_top = Emu(27000)
    cell.margin_bottom = Emu(27000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # 清掉默认空 run
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_CN


def add_quote(slide, text, top=Inches(6.95), color=None, size=12):
    """底部居中金句/结论。"""
    if color is None:
        color = COLOR_PRIMARY
    add_text(slide, Inches(0.6), top, Inches(12.13), Inches(0.45),
             text, size=size, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_bottom_bar(slide, text, top=Inches(6.35), height=Inches(0.55)):
    """底部补充条（白底带边框，左对齐一句话）。"""
    add_rect(slide, Inches(0.6), top, Inches(12.13), height, COLOR_WHITE, COLOR_DIVIDER)
    add_text(slide, Inches(0.7), top, Inches(12), height,
             text, size=11, bold=False, color=COLOR_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def save(prs, out_path):
    out_dir = os.path.dirname(os.path.abspath(out_path))
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path
